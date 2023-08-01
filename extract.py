# We need to import these to fix an AttributeError
import collections
import collections.abc

# Actual imports
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob
import csv

class PowerPointExtractor:
    cur_image_index = 0
    cur_slide_images = []

    invalid_images = []

    input_dir = 'input'
    image_output_dir = 'images'

    def __init__(self, inp, out):
        self.input = inp
        self.image_output_dir = out

    def save_image(self, image, name):
        '''
        Save an image to disk.

        :param image The image to save.
        :param name The partial name of the image, including the directory but not including the sequence number.
        '''

        image_bytes = image.blob

        # Append the sequence number and file extension to the partial name
        name = name + ' ' + str(self.cur_image_index) + '.' + image.ext

        print(name)
        with open(name, 'wb') as f:
            f.write(image_bytes)

        self.cur_image_index += 1
        self.cur_slide_images.append(name.split(os.sep)[1])

    def drill_for_images(self, shape, page, name):
        '''
        Recursive function to look inside grouped shapes for pictures and save them.

        :param shape The parent shape to look inside.
        :param name The partial name for any images we find, including the directory but not including the sequence number.
        '''
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                self.drill_for_images(s, page, name)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                self.save_image(shape.image, name)
            except:
                print(f'Could not process image {shape.name} on slide {page}.')
                self.invalid_images.append(f'Slide {page}: {shape.name}')
                self.cur_slide_images.append(f'INVALID: {shape.name}')

    def generate_image_name_part(self, filename):
        '''
        Generate a partial filename for an image

        :param filename The file name base to use.
        :param output_dir The output directory to use, since we generate the path here too.
        '''
        # Strip input directory path and replace with image output path
        name = filename.split(os.sep)[1]

        # Strip file extension
        name = name.split('.')[0]

        # Prepend output directory to name, for saving later
        return self.image_output_dir + os.sep + name

    def process_files(self):
        """
        Processes a folder of PowerPoint documents and saves a .csv report of all the text within each document, including presenter notes.

        :param input_dir The name of the subfolder to search for PowerPoint documents.
        :param output_dir The name of the subfolder in which to save images.
        """

        if not os.path.exists(self.image_output_dir):
            os.makedirs(self.image_output_dir)

        with open('text.csv', 'w', encoding="utf-8", newline='') as file:
            writer = csv.writer(file)
            presentation_count = 0

            # Write table header
            field = ["File", "Page", "Text", "Notes", "Images"]
            writer.writerow(field)

            # Iterate through files
            print("Processing:")
            for eachfile in glob.glob(self.input_dir + os.sep + "*.pptx"):
                ppt = Presentation(eachfile)
                print("* " + eachfile)
                presentation_count += 1
                self.cur_image_index = 1

                name = self.generate_image_name_part(eachfile)

                # Iterate through slides
                for page, slide in enumerate(ppt.slides):
                    # Collect all the text on the slide into one string, separated by newlines
                    text = ''
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text.strip():
                            text += os.linesep
                            text += shape.text

                    # Collect images from each slide
                    self.cur_slide_images = []

                    # Save images from this slide
                    for shape in slide.shapes:
                        self.drill_for_images(shape, page + 1, name)

                    # Write the page number, collected text, and presenter notes as a new row
                    image_list = ''
                    if (len(self.cur_slide_images) > 0):
                        image_list = self.cur_slide_images
                    writer.writerow([eachfile, page + 1, text, slide.notes_slide.notes_text_frame.text, image_list])

            print("Finished. Total files: " + str(presentation_count))
            if (len(self.invalid_images) > 0):
                print(f'WARNING: {len(self.invalid_images)} invalid images found: {self.invalid_images}')

def main():
    ext = PowerPointExtractor('input', 'images')
    ext.process_files()

if __name__ == "__main__":
    main()
